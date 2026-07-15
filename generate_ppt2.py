"""
Generates FinGenius_AI_Seminar_Presentation.pptx following the exact visual
format of the sample: Federated_IoT_Botnet_Detection.pptx
(dark blue header/footer bars, orange accent strip, white rounded-rectangle
content panels, chained rounded-rectangle + arrow methodology diagram,
styled tables, screenshot panels).
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUTPUT = '/projects/sandbox/summa2/FinGenius_AI_Seminar_Presentation.pptx'
SCREEN_DIR = '/projects/sandbox/summa2'

# ---- Design tokens lifted from the sample deck ----
NAVY = RGBColor.from_string('003092')
NAVY_DARK = RGBColor.from_string('001E5C')
ORANGE = RGBColor.from_string('EF7F1B')
CYAN = RGBColor.from_string('01ADEF')
PANEL_L = RGBColor.from_string('F3F6FB')
PANEL_R = RGBColor.from_string('EEF4FC')
TEXT_DARK = RGBColor.from_string('202433')
TEXT_MUTED = RGBColor.from_string('5A6273')
FOOTER_TEXT_COLOR = RGBColor.from_string('C9D6EE')
WHITE = RGBColor.from_string('FFFFFF')
BODY_TEXT_LIGHT = RGBColor.from_string('E7EEFA')

FONT = 'Calibri'
FOOTER_LINE = 'M.Tech Mini Project  |  [Student Name]  |  [College Name]'

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

prs = Presentation()
prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
prs.slide_height = Emu(int(SLIDE_H_IN * 914400))
blank_layout = prs.slide_layouts[6]

_slide_counter = 0  # used for footer page numbers (title slide excluded)


def _rect(slide, l, t, w, h, fill, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=None, line_spacing=None):
    """runs: list of (text, size_pt, bold, color) OR list of lines, each a list of such tuples."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # Normalize: if runs is a flat list of tuples, treat as single paragraph
    if runs and isinstance(runs[0], tuple):
        lines = [runs]
    else:
        lines = runs

    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for text, size, bold, color in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
    return box



def add_header_footer(slide, eyebrow, title, page_num):
    """Reproduces the sample's header bar (navy block + orange strip + eyebrow
    + title) and footer bar (dark navy strip + credit line + page number)."""
    _rect(slide, 0, 0, SLIDE_W_IN, 1.15, NAVY)
    _rect(slide, 0, 1.15, SLIDE_W_IN, 0.044, ORANGE)
    _textbox(slide, 0.55, 0.13, 9.5, 0.32,
             [(eyebrow, 12.5, True, CYAN)])
    _textbox(slide, 0.55, 0.42, 9.6, 0.68,
             [(title, 25, True, WHITE)])

    _rect(slide, 0, 7.16, SLIDE_W_IN, 0.34, NAVY_DARK)
    _textbox(slide, 0.45, 7.16, 9.5, 0.261,
             [(FOOTER_LINE, 9.5, False, FOOTER_TEXT_COLOR)])
    _textbox(slide, 12.133, 7.16, 0.75, 0.34,
             [(str(page_num), 9.5, True, WHITE)], align=PP_ALIGN.RIGHT)


def add_content_slide(eyebrow, title):
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(blank_layout)
    add_header_footer(slide, eyebrow, title, _slide_counter)
    return slide


def add_panel(slide, l, t, w, h, fill=PANEL_L):
    return slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)).__class__ \
        if False else _rounded_panel(slide, l, t, w, h, fill)


def _rounded_panel(slide, l, t, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    # soften the corner radius similar to sample (~10%)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


def add_bullets(slide, l, t, w, h, items, bullet_color=NAVY, size=13.5):
    """items: list of (bold_lead_text_or_None, rest_text)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    first = True
    for lead, rest in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r0 = p.add_run()
        r0.text = '\u25b8  '
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.name = FONT
        r0.font.color.rgb = bullet_color
        if lead:
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.name = FONT
            r1.font.color.rgb = TEXT_DARK
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.bold = False
        r2.font.name = FONT
        r2.font.color.rgb = TEXT_DARK
    return box



def add_styled_table(slide, l, t, w, h, headers, rows, col_widths_in=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(h))
    table = gtable.table

    if col_widths_in:
        for i, cw in enumerate(col_widths_in):
            table.columns[i].width = Inches(cw)

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.paragraphs[0].text = htext
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(12.5)
        run.font.bold = True
        run.font.name = FONT
        run.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows, start=1):
        band = (r_idx % 2 == 0)
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string('EEF4FC') if band else WHITE
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.paragraphs[0].text = str(val)
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(11.5)
            run.font.name = FONT
            run.font.color.rgb = TEXT_DARK
    return gtable


def add_chain_diagram(slide, top, steps):
    """steps: list of (line1, line2_or_None, fill_color). Draws rounded boxes
    connected by right-pointing arrows, mirroring the sample's methodology row."""
    n = len(steps)
    box_w, box_h, gap_w = 1.62, 0.62, 0.30
    total_w = n * box_w + (n - 1) * gap_w
    left = (SLIDE_W_IN - total_w) / 2
    x = left
    for i, (line1, line2, fill) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
        box.shadow.inherit = False
        try:
            box.adjustments[0] = 0.12
        except Exception:
            pass
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = line1
        r1.font.size = Pt(11.5)
        r1.font.bold = True
        r1.font.name = FONT
        r1.font.color.rgb = WHITE
        if line2:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = line2
            r2.font.size = Pt(9.5)
            r2.font.name = FONT
            r2.font.color.rgb = BODY_TEXT_LIGHT
        x += box_w
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 0.02), Inches(top + 0.18), Inches(gap_w - 0.04), Inches(0.26))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_MUTED
            arrow.line.fill.background()
            arrow.shadow.inherit = False
            x += gap_w



def add_screenshot_panel(slide, l, t, w, h, header_text, header_fill, image_file=None):
    """White rounded panel with a small colored header strip, and either the
    actual screenshot (if available) or a placeholder caption inside."""
    _rounded_panel(slide, l, t, w, h, WHITE)
    header_w = min(w - 0.2, max(2.4, len(header_text) * 0.085))
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l + 0.15), Inches(t + 0.12), Inches(header_w), Inches(0.34))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_fill
    hdr.line.fill.background()
    hdr.shadow.inherit = False
    try:
        hdr.adjustments[0] = 0.25
    except Exception:
        pass
    tf = hdr.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = header_text
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = WHITE

    img_top = t + 0.58
    img_h = h - 0.7
    img_l = l + 0.15
    img_w = w - 0.3

    import os
    if image_file:
        img_path = os.path.join(SCREEN_DIR, image_file)
        if os.path.exists(img_path):
            from PIL import Image
            im = Image.open(img_path)
            ratio = im.width / im.height
            fit_w = img_w
            fit_h = fit_w / ratio
            if fit_h > img_h:
                fit_h = img_h
                fit_w = fit_h * ratio
            pic_l = img_l + (img_w - fit_w) / 2
            pic_t = img_top + (img_h - fit_h) / 2
            slide.shapes.add_picture(img_path, Inches(pic_l), Inches(pic_t), Inches(fit_w), Inches(fit_h))
            return
    # placeholder text if no image supplied
    box = slide.shapes.add_textbox(Inches(img_l), Inches(img_top), Inches(img_w), Inches(img_h))
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = '(Insert Screenshot Here)'
    r2.font.size = Pt(12)
    r2.font.italic = True
    r2.font.name = FONT
    r2.font.color.rgb = TEXT_MUTED



# =====================================================================
# SLIDE 1 - TITLE SLIDE (mirrors sample slide 1 exactly: navy hero block,
# orange strip, white logo card top-left, big title + subtitle, then a
# light "presented by" card below)
# =====================================================================
title_slide = prs.slides.add_slide(blank_layout)
_rect(title_slide, 0, 0, SLIDE_W_IN, 4.55, NAVY)
_rect(title_slide, 0, 4.55, SLIDE_W_IN, 0.056, ORANGE)

_textbox(title_slide, 0.6, 1.95, 12.1, 0.353,
         [('M.TECH MINI PROJECT', 15, True, CYAN)])
_textbox(title_slide, 0.6, 2.35, 12.1, 1.55,
         [[('FinGenius AI - Smart Personal Finance Tracker', 30, True, WHITE)],
          [('An AI Powered MERN Stack Application for Expense, Budget and Goal Management', 30, True, WHITE)]],
         line_spacing=1.05)
_textbox(title_slide, 0.6, 3.78, 12.1, 0.5,
         [('Dashboard: Real-time Financial Overview with Gemini AI Chat Assistant', 16, False, RGBColor.from_string('E6ECF9'))])

_rounded_panel(title_slide, 2.4, 5.05, 8.5, 1.95, PANEL_L)
_textbox(title_slide, 2.75, 5.2, 7.8, 0.4, [('Presented by', 12, True, TEXT_MUTED)])
_textbox(title_slide, 2.75, 5.52, 7.8, 0.42, [('[Student Name]  |  Reg. No. [Register Number]', 19, True, NAVY)])
_textbox(title_slide, 2.75, 5.98, 7.8, 0.337, [('M.Tech, Computer Science Engineering', 14, False, TEXT_DARK)])
_textbox(title_slide, 2.75, 6.3, 7.8, 0.32, [('Department of Computer Science and Engineering', 13, False, TEXT_MUTED)])
_textbox(title_slide, 2.75, 6.6, 7.8, 0.34, [('[College Name]', 13, False, TEXT_MUTED)])



# =====================================================================
# SLIDE 2 - MOTIVATION: Problem Statement & Objectives (two side-by-side panels)
# =====================================================================
s2 = add_content_slide('MOTIVATION', 'Problem Statement & Objectives')
_rounded_panel(s2, 0.55, 1.45, 6.15, 5.35, PANEL_L)
_textbox(s2, 0.85, 1.62, 5.6, 0.4, [('The Problem', 17, True, NAVY)])
add_bullets(s2, 0.85, 2.1, 5.6, 4.5, [
    ('Scattered records: ', 'people track expenses in notebooks or spreadsheets, with no single place for income, budget, and goals.'),
    ('No intelligent insight: ', 'free tools simply store numbers without explaining spending patterns or suggesting improvements.'),
    ('Manual entry is tedious: ', 'typing every transaction by hand causes many users to abandon tracking within weeks.'),
    ('No natural language help: ', 'users cannot simply ask "How much did I spend on food?" and get a direct answer.'),
    ('Privacy concerns: ', 'many finance apps need access to bank SMS or account details to work.'),
], bullet_color=NAVY, size=13)

_rounded_panel(s2, 6.9, 1.45, 5.9, 5.35, PANEL_R)
_textbox(s2, 7.2, 1.62, 5.3, 0.4, [('Project Objectives', 17, True, NAVY)])
add_bullets(s2, 7.2, 2.1, 5.35, 4.5, [
    ('Unified tracking: ', 'record expenses and income and organise them into clear categories.'),
    ('Budget and goal planning: ', 'set monthly category budgets and savings goals with progress tracking.'),
    ('AI powered chat: ', 'integrate Google Gemini AI to answer questions using the user\'s own data.'),
    ('Faster data entry: ', 'support receipt scanning (OCR) and voice based expense entry.'),
    ('Visual analytics: ', 'present spending and saving trends through interactive charts.'),
], bullet_color=CYAN, size=13)



# =====================================================================
# SLIDE 3 - SYSTEM ARCHITECTURE (two labelled phase panels, mirrors sample slide 3)
# =====================================================================
s3 = add_content_slide('SYSTEM DESIGN', 'System Architecture - Frontend & Backend')

hdr1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.45), Inches(6.05), Inches(0.42))
hdr1.fill.solid(); hdr1.fill.fore_color.rgb = NAVY; hdr1.line.fill.background(); hdr1.shadow.inherit = False
try:
    hdr1.adjustments[0] = 0.2
except Exception:
    pass
tf = hdr1.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = 'CLIENT SIDE - React.js Frontend'
r.font.size = Pt(13); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE

_rounded_panel(s3, 0.5, 1.95, 6.05, 5.05, WHITE)
add_bullets(s3, 0.85, 2.25, 5.4, 4.6, [
    ('React 18 + Vite: ', 'component based single page application'),
    ('Tailwind CSS: ', 'responsive glassmorphism styling, dark/light theme'),
    ('Framer Motion: ', 'smooth page and card animations'),
    ('Recharts: ', 'interactive analytics charts'),
    ('React Router v6: ', 'client side routing between pages'),
    ('Axios: ', 'REST API communication with the backend'),
    ('Pages: ', 'Landing, Login, Register, Dashboard, Expenses, Income, Budget, Goals, Analytics, Chat, Settings'),
], bullet_color=NAVY, size=12.5)

hdr2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.45), Inches(6.05), Inches(0.42))
hdr2.fill.solid(); hdr2.fill.fore_color.rgb = ORANGE; hdr2.line.fill.background(); hdr2.shadow.inherit = False
try:
    hdr2.adjustments[0] = 0.2
except Exception:
    pass
tf = hdr2.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = 'SERVER SIDE - Node.js Backend'
r.font.size = Pt(13); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE

_rounded_panel(s3, 6.75, 1.95, 6.05, 5.05, WHITE)
add_bullets(s3, 7.1, 2.25, 5.4, 4.6, [
    ('Node.js + Express.js: ', 'REST API server with route-controller structure'),
    ('MongoDB Atlas + Mongoose: ', 'cloud database and schema modelling'),
    ('JWT + bcryptjs: ', 'token based authentication, hashed passwords'),
    ('Google Gemini AI: ', 'AI chat assistant with financial context'),
    ('Tesseract.js: ', 'OCR engine for receipt scanning'),
    ('Multer: ', 'file upload handling for receipt images'),
    ('Middleware: ', 'JWT auth guard, rate limiter, input validator, error handler'),
], bullet_color=ORANGE, size=12.5)



# =====================================================================
# SLIDE 4 - HOW IT WORKS: Methodology (chain diagram + explanation bullets,
# mirrors the sample's 7-box methodology chain exactly in style)
# =====================================================================
s4 = add_content_slide('HOW IT WORKS', 'Methodology - From Login to AI Insight')

add_chain_diagram(s4, 1.42, [
    ('Register &\nLogin', 'JWT + bcrypt', NAVY),
    ('Dashboard', 'Balance overview', CYAN),
    ('Add\nTransaction', 'Manual/OCR/Voice', NAVY),
    ('Budget &\nGoals', 'Category limits', ORANGE),
    ('Analytics', 'Recharts trends', NAVY),
    ('AI Chat', 'Gemini financial Q&A', CYAN),
])

_textbox(s4, 0.55, 2.35, 12.3, 4.4, [
    [('User Authentication: ', 14, True, NAVY), ('a user registers with name, email, and password; passwords are hashed with bcryptjs (12 salt rounds) and login issues a 7-day JWT token used to authorise every later request.', 14, False, TEXT_DARK)],
    [('Recording Transactions: ', 14, True, NAVY), ('expenses and income can be entered manually, extracted from a receipt photo using Tesseract.js OCR, or logged by speaking through the Web Speech API.', 14, False, TEXT_DARK)],
    [('Budget and Goal Tracking: ', 14, True, NAVY), ('users set category wise monthly budgets and savings goals with a target amount; the system tracks actual spending and contribution progress automatically.', 14, False, TEXT_DARK)],
    [('Analytics Generation: ', 14, True, NAVY), ('Recharts renders monthly expense trend, income vs expense, category breakdown, and savings trend charts from the stored MongoDB data.', 14, False, TEXT_DARK)],
    [('AI Powered Chat: ', 14, True, NAVY), ('the backend sends the user\u2019s financial summary (expenses by category, income, budget status, goals, recent transactions) to Google Gemini 1.5 Flash, which answers questions in plain language.', 14, False, TEXT_DARK)],
], line_spacing=1.35)



# =====================================================================
# SLIDE 5 - TECHNOLOGY STACK (table on the left, screenshot panel on the
# right + observations, mirrors the sample's "Results" table+chart layout)
# =====================================================================
s5 = add_content_slide('IMPLEMENTATION', 'Technology Stack')

_textbox(s5, 0.55, 1.5, 6.0, 0.35, [('Frontend and Backend Technologies', 14, True, TEXT_DARK)])
add_styled_table(s5, 0.55, 1.9, 6.3, 3.4,
                  ['Technology', 'Purpose'],
                  [
                      ['React 18 + Vite', 'Frontend UI'],
                      ['Tailwind CSS', 'Responsive styling'],
                      ['Node.js + Express.js', 'Backend REST APIs'],
                      ['MongoDB + Mongoose', 'Database & schema'],
                      ['JWT + bcryptjs', 'Auth & security'],
                      ['Google Gemini AI', 'AI chat assistant'],
                      ['Tesseract.js', 'Receipt OCR'],
                      ['Web Speech API', 'Voice entry'],
                      ['Recharts', 'Analytics charts'],
                  ],
                  col_widths_in=[3.0, 3.3])

add_screenshot_panel(s5, 7.1, 1.45, 5.7, 5.05, 'Dashboard - Financial Overview', NAVY, 'Screenshot 2026-07-13 205528.png')



# =====================================================================
# SLIDE 6 - LIVE SYSTEM: Application Walkthrough (two screenshot panels,
# mirrors the sample's "Dashboard Demonstration" slide)
# =====================================================================
s6 = add_content_slide('LIVE SYSTEM', 'Application Walkthrough')

add_screenshot_panel(s6, 0.5, 1.42, 5.9, 5.05, 'Expenses \u2022 Search \u2022 Filter \u2022 Scan \u2022 Voice', NAVY, 'Screenshot 2026-07-13 205600.png')
add_screenshot_panel(s6, 6.9, 1.42, 5.9, 5.05, 'AI Chat \u2022 Gemini Powered Assistant', ORANGE, 'Screenshot 2026-07-13 205745.png')



# =====================================================================
# SLIDE 7 - NOVELTY & CONTRIBUTIONS: Comparison with Existing Systems
# (table + Merits/Limitations two panels, mirrors sample slide 7)
# =====================================================================
s7 = add_content_slide('COMPARISON', 'Comparison with Existing Systems')

_textbox(s7, 0.55, 1.5, 8.0, 0.35, [('Existing Methods vs. FinGenius AI', 14, True, TEXT_DARK)])
add_styled_table(s7, 0.55, 1.9, 12.25, 2.1,
                  ['Feature', 'Manual / Excel', 'Typical Finance Apps', 'FinGenius AI'],
                  [
                      ['Security', 'None / basic', 'Varies by app', 'JWT + bcrypt hashing'],
                      ['AI Assistance', 'Not available', 'Rare or basic', 'Gemini AI chatbot'],
                      ['Expense Entry', 'Manual only', 'Manual entry', 'Manual, OCR, and voice'],
                      ['Cost', 'Free', 'Free to paid', 'Free for this project'],
                  ],
                  col_widths_in=[2.6, 2.9, 3.35, 3.4])

hdr_m = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.35), Inches(6.05), Inches(0.4))
hdr_m.fill.solid(); hdr_m.fill.fore_color.rgb = NAVY; hdr_m.line.fill.background(); hdr_m.shadow.inherit = False
try:
    hdr_m.adjustments[0] = 0.2
except Exception:
    pass
tfm = hdr_m.text_frame; tfm.paragraphs[0].alignment = PP_ALIGN.CENTER
rm = tfm.paragraphs[0].add_run(); rm.text = 'Merits'
rm.font.size = Pt(13); rm.font.bold = True; rm.font.name = FONT; rm.font.color.rgb = WHITE

add_bullets(s7, 0.55, 4.9, 6.05, 1.9, [
    (None, 'AI powered insights based on the user\u2019s real financial data'),
    (None, 'Multiple ways to log expenses - manual, receipt scan, voice'),
    (None, 'Completely free with all features unlocked'),
    (None, 'Secure - JWT authentication and hashed passwords'),
], bullet_color=NAVY, size=12)

hdr_l = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(4.35), Inches(5.95), Inches(0.4))
hdr_l.fill.solid(); hdr_l.fill.fore_color.rgb = ORANGE; hdr_l.line.fill.background(); hdr_l.shadow.inherit = False
try:
    hdr_l.adjustments[0] = 0.2
except Exception:
    pass
tfl = hdr_l.text_frame; tfl.paragraphs[0].alignment = PP_ALIGN.CENTER
rl = tfl.paragraphs[0].add_run(); rl.text = 'Limitations'
rl.font.size = Pt(13); rl.font.bold = True; rl.font.name = FONT; rl.font.color.rgb = WHITE

add_bullets(s7, 6.85, 4.9, 5.95, 1.9, [
    (None, 'No direct bank account integration yet - entries are manual or scanned'),
    (None, 'Requires an internet connection for all operations'),
    (None, 'OCR accuracy depends on the quality of the receipt image'),
    (None, 'Voice recognition works best in Chrome or Edge browsers'),
], bullet_color=ORANGE, size=12)



# =====================================================================
# SLIDE 8 - WRAP-UP: Conclusion & Future Work (mirrors sample's final slide)
# =====================================================================
s8 = add_content_slide('WRAP-UP', 'Conclusion & Future Work')

hdr_c = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.42), Inches(6.05), Inches(0.4))
hdr_c.fill.solid(); hdr_c.fill.fore_color.rgb = NAVY; hdr_c.line.fill.background(); hdr_c.shadow.inherit = False
try:
    hdr_c.adjustments[0] = 0.2
except Exception:
    pass
tfc = hdr_c.text_frame; tfc.paragraphs[0].alignment = PP_ALIGN.CENTER
rc = tfc.paragraphs[0].add_run(); rc.text = 'Conclusion'
rc.font.size = Pt(13); rc.font.bold = True; rc.font.name = FONT; rc.font.color.rgb = WHITE

add_bullets(s8, 0.55, 1.95, 6.05, 4.6, [
    (None, 'Delivered a full-stack personal finance tracker combining expense, income, budget, and goal management with an AI chat assistant.'),
    (None, 'Integrated Google Gemini AI so users get personalised, data based answers instead of generic financial tips.'),
    (None, 'Added Receipt OCR and Voice Entry to reduce the effort of manual data entry.'),
    (None, 'Implemented secure JWT authentication and a fully responsive, modern user interface.'),
    (None, 'Provides a practical, free, all-in-one alternative to scattered manual and spreadsheet based tracking.'),
], bullet_color=NAVY, size=13)

hdr_f = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.42), Inches(5.95), Inches(0.4))
hdr_f.fill.solid(); hdr_f.fill.fore_color.rgb = ORANGE; hdr_f.line.fill.background(); hdr_f.shadow.inherit = False
try:
    hdr_f.adjustments[0] = 0.2
except Exception:
    pass
tff = hdr_f.text_frame; tff.paragraphs[0].alignment = PP_ALIGN.CENTER
rf = tff.paragraphs[0].add_run(); rf.text = 'Future Work'
rf.font.size = Pt(13); rf.font.bold = True; rf.font.name = FONT; rf.font.color.rgb = WHITE

add_bullets(s8, 6.85, 1.95, 5.95, 4.6, [
    (None, 'Bank account and UPI integration for automatic transaction import'),
    (None, 'Improving OCR accuracy for a wider variety of receipt formats'),
    (None, 'Multi-language support for the UI and voice entry'),
    (None, 'Investment tracking alongside income and expenses'),
    (None, 'Predictive expense forecasting using past spending patterns'),
    (None, 'A dedicated mobile application with push notifications'),
], bullet_color=ORANGE, size=13)

_rect(s8, 0, 6.55, SLIDE_W_IN, 0.61, NAVY)
_textbox(s8, 0, 6.55, SLIDE_W_IN, 0.61, [('Thank You', 20, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUTPUT)
print(f'Presentation saved to {OUTPUT}')
