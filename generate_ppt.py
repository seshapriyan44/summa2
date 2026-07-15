"""
Generates a seminar/viva presentation for FinGenius AI, following the exact
slide format (theme, layouts, title slide style, footer/date/slide-number
placeholders, bullet + sub-bullet content, table slide, block-diagram slide)
of the sample PPT:
'127033016 Integration of cloud manufacturing and cyber physical systems.pptx'
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

TEMPLATE = '/projects/sandbox/summa2/127033016 Integration of cloud manufacturing and cyber physical systems.pptx'
OUTPUT = '/projects/sandbox/summa2/FinGenius_AI_Seminar_Presentation.pptx'

FOOTER_TEXT = 'FinGenius AI - Smart Personal Finance Tracker'

prs = Presentation(TEMPLATE)

# Keep master + layouts + theme, remove all existing (sample) slides.
# We must also drop the relationship (and therefore the underlying part)
# for each slide, otherwise python-pptx re-adds new slides using the same
# part names, producing duplicate-name warnings and a corrupt package.
xml_slides = prs.slides._sldIdLst
rids_to_drop = []
for sld in list(xml_slides):
    rids_to_drop.append(sld.get(qn('r:id')))
    xml_slides.remove(sld)
for rid in rids_to_drop:
    prs.part.drop_rel(rid)

title_layout = prs.slide_masters[0].slide_layouts[0]   # Title Slide
content_layout = prs.slide_masters[0].slide_layouts[1]  # Title and Content


def clone_placeholder(slide, layout, idx):
    """Copy a placeholder shape (e.g. date/footer/slide number) from the
    layout into the slide, since add_slide() only copies title/body/etc."""
    src = None
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == idx:
            src = ph
            break
    if src is None:
        return None
    sp = copy.deepcopy(src.element)
    slide.shapes._spTree.append(sp)
    return slide.placeholders[idx]


def add_footer_date_slidenum(slide, layout):
    clone_placeholder(slide, layout, 10)  # date
    ftr = clone_placeholder(slide, layout, 11)  # footer
    clone_placeholder(slide, layout, 12)  # slide number
    if ftr is not None:
        ftr.text_frame.text = FOOTER_TEXT


def add_title_slide():
    slide = prs.slides.add_slide(title_layout)
    subtitle = slide.placeholders[1]
    tf = subtitle.text_frame
    tf.clear()

    lines = [
        ('FinGenius AI - Smart Personal Finance Tracker', True),
        ('An AI Powered MERN Stack Application for Expense, Budget and Goal Management', True),
        ('Name:', True, ' [Student Name]', False),
        ('Regr No.:', True, ' [Register Number]', False),
        ('M.Tech. CSE', False),
    ]

    first = True
    for item in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if len(item) == 2:
            text, bold = item
            run = p.add_run()
            run.text = text
            run.font.bold = bold
        else:
            text1, bold1, text2, bold2 = item
            r1 = p.add_run()
            r1.text = text1
            r1.font.bold = bold1
            r2 = p.add_run()
            r2.text = text2
            r2.font.bold = bold2

    add_footer_date_slidenum(slide, title_layout)
    return slide


def add_bullet_slide(title_text, bullets):
    """bullets: list of (level, text) tuples, level 0 or 1"""
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text_frame.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.text = text
    add_footer_date_slidenum(slide, content_layout)
    return slide


def add_table_slide(title_text, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text_frame.text = title_text
    # remove the empty content placeholder body (we'll place a table instead)
    body_ph = slide.placeholders[1]
    body_ph._element.getparent().remove(body_ph._element)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(12.3)
    height = Inches(0.6 * n_rows)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gtable.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)

    add_footer_date_slidenum(slide, content_layout)
    return slide



def add_blank_title_slide(title_text):
    """A slide with just a title placeholder (like the sample's diagram slides),
    onto which we can freely place shapes."""
    slide = prs.slides.add_slide(content_layout)
    slide.placeholders[0].text_frame.text = title_text
    body_ph = slide.placeholders[1]
    body_ph._element.getparent().remove(body_ph._element)
    add_footer_date_slidenum(slide, content_layout)
    return slide


def add_box(slide, text, left, top, width, height, fill_hex='5B9BD5', font_size=12, font_color='FFFFFF'):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))  # 1 = RECTANGLE
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    box.line.color.rgb = RGBColor.from_string('000000')
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.color.rgb = RGBColor.from_string(font_color)
            r.font.bold = True
    return box


def add_down_arrow(slide, left, top, height=0.3):
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(left), Inches(top), Inches(left), Inches(top + height)
    )
    connector.line.color.rgb = RGBColor.from_string('000000')
    connector.line.width = Pt(2)
    return connector



# =====================================================================
# SLIDE 1 - TITLE SLIDE
# =====================================================================
add_title_slide()

# =====================================================================
# SLIDE 2 - PRESENTATION OUTLINE
# =====================================================================
add_bullet_slide('Presentation Outline', [
    (0, 'Background and motivation'),
    (0, 'Problem statement and objectives'),
    (0, 'Existing personal finance tracking methods'),
    (0, 'Overview of the proposed system - FinGenius AI'),
    (0, 'System architecture and workflow'),
    (0, 'Module description and database design'),
    (0, 'Advantages and limitations of the approach'),
    (0, 'Conclusion and future scope'),
])

# =====================================================================
# SLIDE 3 - PROBLEM STATEMENT
# =====================================================================
add_bullet_slide('Problem Statement', [
    (0, 'Most people find it difficult to track daily expenses, income, and savings in one place.'),
    (0, 'Free tools rarely provide intelligent, personalised financial advice.'),
    (0, 'A web based application is proposed that combines expense, income, budget, and goal '
        'tracking with an AI chat assistant, receipt scanning, and voice entry.'),
])

# =====================================================================
# SLIDE 4 - INTRODUCTION
# =====================================================================
add_bullet_slide('Introduction', [
    (0, 'Managing personal finance is a challenge faced by most individuals today'),
    (0, 'Traditional methods (notebooks, spreadsheets) are manual and offer no insights'),
    (0, 'FinGenius AI is a smart personal finance tracker built on the MERN stack'),
    (0, 'Integrates Google Gemini AI to answer questions about the user\'s own data'),
    (1, 'Acts like a personal financial advisor'),
    (0, 'Supports receipt scanning (OCR) and voice based expense entry'),
    (0, 'Provides dashboards, budgets, goals, and analytics in a single platform'),
])



# =====================================================================
# SLIDE 5 - EXISTING METHODS (table, mirrors "Literature Survey" table slide)
# =====================================================================
add_table_slide(
    'Existing Methods for Personal Finance Tracking',
    ['S. No.', 'Method', 'Description', 'Limitation'],
    [
        ['1.', 'Manual / Notebook', 'Writing income and expenses by hand', 'Slow, error-prone, no analysis'],
        ['2.', 'Spreadsheet (Excel)', 'Recording transactions with manual formulas', 'Needs technical effort, no AI insight'],
        ['3.', 'Mobile Banking Apps', 'Shows transactions of one bank account', 'No categorisation, no cash tracking'],
        ['4.', 'Third-Party Finance Apps', 'Apps like Money Manager, Walnut, Mint', 'Advanced features often paid, privacy concerns'],
    ],
    col_widths=[1.0, 2.8, 4.5, 4.0],
)

# =====================================================================
# SLIDE 6 - PROBLEMS IN EXISTING SYSTEMS
# =====================================================================
add_bullet_slide('Problems in Existing Systems', [
    (0, 'Limited intelligence - most tools only record data, without suggestions'),
    (0, 'No AI chatbot to answer natural language questions about spending'),
    (0, 'No receipt scanning support in most free versions'),
    (0, 'Very few apps support voice based expense entry'),
    (0, 'Privacy concerns - many apps need access to bank SMS or accounts'),
    (0, 'Useful features like analytics and budget alerts are often paid'),
])



# =====================================================================
# SLIDE 7 - OVERVIEW OF PROPOSED SYSTEM
# =====================================================================
add_bullet_slide('Overview of Proposed System', [
    (0, 'FinGenius AI - an AI powered personal finance tracker'),
    (0, 'Built using the MERN stack (MongoDB, Express.js, React, Node.js)'),
    (0, 'Client-server architecture with RESTful APIs'),
    (0, 'Google Gemini AI integrated as a chat based financial assistant'),
    (0, 'Modern glassmorphism UI with dark and light theme support'),
    (0, 'Fully responsive - works on desktop, tablet, and mobile'),
])

# =====================================================================
# SLIDE 8 - SYSTEM ARCHITECTURE (Text layers, similar style to "Architecture" list slide)
# =====================================================================
add_bullet_slide('System Architecture', [
    (0, 'Presentation Layer - React.js with Tailwind CSS'),
    (1, 'Dashboard, Expenses, Income, Budget, Goals, Analytics, AI Chat'),
    (0, 'Business Logic Layer - Node.js with Express.js'),
    (1, 'Auth, Expense, Income, Budget, Goal, Analytics, Chat, OCR controllers'),
    (0, 'Middleware Layer - JWT authentication, rate limiter, validator'),
    (0, 'Data and AI Layer - MongoDB Atlas, Google Gemini AI, Tesseract.js OCR'),
])

# =====================================================================
# SLIDE 9 - ARCHITECTURE DIAGRAM (block diagram, mirrors sample's picture layer slide)
# =====================================================================
slide9 = add_blank_title_slide('Architecture Diagram')
add_box(slide9, 'Presentation Layer\n(React + Tailwind CSS)', 3.9, 1.5, 5.5, 0.9, '5B9BD5')
add_down_arrow(slide9, 6.65, 2.4, 0.3)
add_box(slide9, 'Business Logic Layer\n(Node.js + Express.js)', 3.9, 2.7, 5.5, 0.9, 'ED7D31')
add_down_arrow(slide9, 6.65, 3.6, 0.3)
add_box(slide9, 'Middleware\n(JWT Auth, Rate Limiter, Validator)', 3.9, 3.9, 5.5, 0.7, 'A5A5A5')
add_down_arrow(slide9, 6.65, 4.6, 0.3)
add_box(slide9, 'MongoDB Atlas\n(Database)', 1.8, 4.9, 2.9, 0.9, '4472C4')
add_box(slide9, 'Google Gemini AI\n(Chat Assistant)', 5.15, 4.9, 2.9, 0.9, '70AD47')
add_box(slide9, 'Tesseract.js\n(OCR Engine)', 8.5, 4.9, 2.9, 0.9, 'FFC000')

# =====================================================================
# SLIDE 10 - WORKFLOW
# =====================================================================
add_bullet_slide('Workflow', [
    (0, 'User registers and logs in using a JWT secured account'),
    (0, 'Dashboard shows balance, income, expense, and savings summary'),
    (0, 'User adds transactions - manually, via receipt scan, or by voice'),
    (0, 'User creates monthly budgets and sets financial goals'),
    (0, 'Analytics module shows trends using interactive charts'),
    (0, 'User can chat with the Gemini AI assistant about their own finances'),
])



# =====================================================================
# SLIDE 11 - AUTHENTICATION MODULE
# =====================================================================
add_bullet_slide('Authentication Module', [
    (0, 'Registration with name, email, and password'),
    (0, 'Passwords hashed using bcryptjs with 12 salt rounds'),
    (0, 'Login generates a JWT token valid for 7 days'),
    (1, 'Token used to authorise all further API requests'),
    (0, 'Forgot password and reset password via email link'),
    (0, 'Protected routes redirect unauthenticated users to login'),
])

# =====================================================================
# SLIDE 12 - DASHBOARD MODULE
# =====================================================================
add_bullet_slide('Dashboard Module', [
    (0, 'Shows Total Balance, Monthly Income, Monthly Expense, Monthly Savings'),
    (1, 'Each card also shows percentage change from previous month'),
    (0, 'Quick action buttons - Add Expense, Add Income, AI Chat, Scan Receipt'),
    (0, 'Income vs Expense chart and Expense by Category chart'),
    (0, 'Personalised greeting based on time of day'),
])

# =====================================================================
# SLIDE 13 - EXPENSE AND INCOME MANAGEMENT
# =====================================================================
add_bullet_slide('Expense and Income Management', [
    (0, 'Add, edit, and delete expenses with category, amount, date, payment method'),
    (0, 'Add income entries with source, amount, date, and recurring option'),
    (1, 'Sources: Salary, Freelancing, Investments, Business, Rental, Other'),
    (0, 'Search and filter transactions by description or category'),
    (0, 'Export option to download transaction data'),
])

# =====================================================================
# SLIDE 14 - BUDGET AND GOAL PLANNING
# =====================================================================
add_bullet_slide('Budget and Goal Planning', [
    (0, 'Budget module - set monthly spending limits per category'),
    (1, 'Progress bars track actual spending against the set limit'),
    (0, 'Goal module - create savings goals with a target amount and date'),
    (1, 'Users can contribute towards a goal and track progress'),
    (1, 'Goals organised as Active, Achieved, Paused, and All'),
])



# =====================================================================
# SLIDE 15 - ANALYTICS MODULE
# =====================================================================
add_bullet_slide('Analytics Module', [
    (0, 'Monthly Expense Trend - last 12 months as a line chart'),
    (0, 'Income vs Expense - comparison chart over 12 months'),
    (0, 'Expense by Category - pie / donut chart for the current month'),
    (0, 'Savings Trend - shows income minus expense over time'),
    (0, 'Charts built using the Recharts library, update in real time'),
])

# =====================================================================
# SLIDE 16 - AI CHAT MODULE (GEMINI AI)
# =====================================================================
add_bullet_slide('AI Chat Module (Gemini AI)', [
    (0, 'Google Gemini 1.5 Flash used as a personal financial advisor'),
    (0, 'Backend sends the user\'s financial summary along with the question'),
    (1, 'Expenses by category, income, budget status, active goals, recent transactions'),
    (0, 'Users can ask - "How much did I spend this month?", "Can I afford a laptop?"'),
    (0, 'Chat history is saved and shown in a sidebar for later reference'),
])

# =====================================================================
# SLIDE 17 - RECEIPT OCR AND VOICE ENTRY
# =====================================================================
add_bullet_slide('Receipt OCR and Voice Entry', [
    (0, 'Receipt OCR - uses Tesseract.js to read text from a receipt image'),
    (1, 'Extracts merchant name, amount, date, and suggests a category'),
    (0, 'Voice Entry - uses the browser\'s Web Speech API'),
    (1, 'User can say "I spent 250 rupees on petrol" to log an expense'),
    (0, 'Both features reduce the effort of manual data entry'),
])

# =====================================================================
# SLIDE 18 - DATABASE DESIGN
# =====================================================================
add_bullet_slide('Database Design (MongoDB)', [
    (0, 'MongoDB Atlas used as the cloud database'),
    (0, 'User collection - stores name, email, hashed password'),
    (0, 'Expense and Income collections - linked to a user by userId'),
    (0, 'Budget collection - monthly limit per category'),
    (0, 'Goal collection - target amount, current amount, contributions'),
    (0, 'Chat and Notification collections - AI conversations and alerts'),
])

# =====================================================================
# SLIDE 19 - TECHNOLOGY STACK (table)
# =====================================================================
add_table_slide(
    'Technology Stack',
    ['Technology', 'Purpose'],
    [
        ['React.js + Vite', 'Frontend user interface'],
        ['Tailwind CSS', 'Responsive glassmorphism styling'],
        ['Node.js + Express.js', 'Backend server and REST APIs'],
        ['MongoDB + Mongoose', 'Database and schema modelling'],
        ['JWT + bcryptjs', 'Authentication and password security'],
        ['Google Gemini AI', 'AI chat assistant'],
        ['Tesseract.js', 'Receipt OCR'],
        ['Web Speech API', 'Voice based expense entry'],
        ['Recharts', 'Analytics charts'],
    ],
    col_widths=[4.5, 7.8],
)



# =====================================================================
# SLIDE 20 - ADVANTAGES OF PROPOSED APPROACH
# =====================================================================
add_bullet_slide('Advantages of Proposed Approach', [
    (0, 'AI powered insights based on the user\'s real financial data'),
    (0, 'Multiple ways to log expenses - manual, receipt scan, voice'),
    (0, 'Completely free with all features unlocked'),
    (0, 'Modern, responsive user interface with dark and light mode'),
    (0, 'Secure - JWT authentication and hashed passwords'),
    (0, 'Combines expense, income, budget, goal, and analytics in one app'),
])

# =====================================================================
# SLIDE 21 - CHALLENGES AND LIMITATIONS
# =====================================================================
add_bullet_slide('Challenges and Limitations', [
    (0, 'No direct bank account integration - entries are manual or scanned'),
    (0, 'Requires an internet connection for all operations'),
    (0, 'OCR accuracy depends on image quality of the receipt'),
    (0, 'Voice recognition works best in Chrome or Edge browsers'),
    (0, 'Designed for a single user, not for shared family accounts'),
    (0, 'AI chat feature depends on the Gemini API being available'),
])

# =====================================================================
# SLIDE 22 - COMPARISON WITH EXISTING SYSTEM (table)
# =====================================================================
add_table_slide(
    'Comparison with Existing System',
    ['Feature', 'Manual / Excel', 'Typical Finance Apps', 'FinGenius AI'],
    [
        ['Security', 'None / basic', 'Varies by app', 'JWT + bcrypt hashing'],
        ['AI Assistance', 'Not available', 'Rare or basic', 'Gemini AI chatbot'],
        ['Expense Entry', 'Manual only', 'Manual entry', 'Manual, OCR, and voice'],
        ['Budget Planning', 'Manual formulas', 'Basic, often paid', 'Category wise with alerts'],
        ['Analytics', 'Basic charts', 'Standard charts', 'Interactive Recharts'],
        ['Cost', 'Free', 'Free to paid', 'Free for this project'],
    ],
    col_widths=[3.0, 3.0, 3.3, 4.0],
)

# =====================================================================
# SLIDE 23 - FUTURE WORK
# =====================================================================
add_bullet_slide('Future Work', [
    (0, 'Bank account and UPI integration for automatic transaction import'),
    (0, 'Improving OCR accuracy for a wider variety of receipt formats'),
    (0, 'Multi-language support for UI and voice entry'),
    (0, 'Investment tracking alongside income and expenses'),
    (0, 'Predictive expense forecasting using past spending patterns'),
    (0, 'Dedicated mobile application with push notifications'),
])

# =====================================================================
# SLIDE 24 - CONCLUSION
# =====================================================================
add_bullet_slide('Conclusion', [
    (0, 'FinGenius AI combines the MERN stack with Gemini AI for smarter finance tracking'),
    (0, 'Covers expenses, income, budgets, goals, and analytics in one platform'),
    (0, 'Receipt OCR and voice entry reduce the effort of manual data entry'),
    (0, 'AI chat assistant gives personalised, data based financial answers'),
    (0, 'Provides a strong, extensible foundation for future enhancement'),
])

prs.save(OUTPUT)
print(f'Presentation saved to {OUTPUT}')
