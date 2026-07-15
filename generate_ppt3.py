"""
Generates FinGenius_AI_Seminar_Presentation.pptx using a DARK GLASSMORPHISM
theme that matches the actual FinGenius AI product branding (indigo / green /
amber accents on a deep navy-black background with translucent "glass" cards),
while keeping the same overall slide FORMAT used previously: a fixed
header/footer chrome on every slide, side-by-side panels, a chained
step-diagram, styled tables, and screenshot panels.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = '/projects/sandbox/summa2/FinGenius_AI_Seminar_Presentation.pptx'
SCREEN_DIR = '/projects/sandbox/summa2'

# ---- Design tokens: FinGenius AI's own product palette (dark glass UI) ----
BG_DARK = RGBColor.from_string('0F1120')       # deep space background
BG_PANEL = RGBColor.from_string('1A1D33')      # glass card background
BG_PANEL_ALT = RGBColor.from_string('20233F')  # alternate glass card
INDIGO = RGBColor.from_string('6366F1')        # primary brand color
GREEN = RGBColor.from_string('22C55E')         # MERN / success accent
AMBER = RGBColor.from_string('F59E0B')         # AI accent
BORDER = RGBColor.from_string('353A5C')        # subtle card border
TEXT_LIGHT = RGBColor.from_string('F1F3FF')    # primary text on dark
TEXT_MUTED = RGBColor.from_string('9CA3C7')    # secondary text on dark
WHITE = RGBColor.from_string('FFFFFF')

FONT = 'Calibri'
FOOTER_LINE = 'FinGenius AI  \u2022  M.Tech Mini Project  \u2022  [Student Name]  \u2022  [College Name]'

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

prs = Presentation()
prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
prs.slide_height = Emu(int(SLIDE_H_IN * 914400))
blank_layout = prs.slide_layouts[6]

_slide_counter = 0


def _rect(slide, l, t, w, h, fill, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w or 0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _rounded(slide, l, t, w, h, fill, radius=0.06, line_color=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def _textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=None, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if runs and isinstance(runs[0], tuple):
        lines = [runs]
    else:
        lines = runs

    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for text, size, bold, color in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
    return box



def add_header_footer(slide, eyebrow, title, page_num, accent=INDIGO):
    """Dark chrome: full dark background, a slim accent-colored top rule,
    eyebrow label + title in the top band, and a footer strip with credit
    line + page number - same FORMAT as before but restyled dark/glass."""
    _rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, BG_DARK)
    _rect(slide, 0, 0, SLIDE_W_IN, 0.09, accent)
    _textbox(slide, 0.55, 0.35, 9.5, 0.32, [(eyebrow, 12.5, True, accent)])
    _textbox(slide, 0.55, 0.66, 9.6, 0.6, [(title, 24, True, TEXT_LIGHT)])
    _rect(slide, 0, 1.38, SLIDE_W_IN, 0.02, BORDER)

    _rect(slide, 0, 7.16, SLIDE_W_IN, 0.34, RGBColor.from_string('090A16'))
    _textbox(slide, 0.45, 7.16, 9.5, 0.261, [(FOOTER_LINE, 9.5, False, TEXT_MUTED)])
    _textbox(slide, 12.133, 7.16, 0.75, 0.34, [(str(page_num), 9.5, True, accent)], align=PP_ALIGN.RIGHT)


def add_content_slide(eyebrow, title, accent=INDIGO):
    global _slide_counter
    _slide_counter += 1
    slide = prs.slides.add_slide(blank_layout)
    add_header_footer(slide, eyebrow, title, _slide_counter, accent)
    return slide


def add_glass_panel(slide, l, t, w, h, fill=BG_PANEL, border=BORDER):
    return _rounded(slide, l, t, w, h, fill, radius=0.05, line_color=border, line_w=1.0)


def add_bullets(slide, l, t, w, h, items, bullet_color=INDIGO, size=13.5, text_color=TEXT_LIGHT):
    """items: list of (bold_lead_text_or_None, rest_text)"""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    first = True
    for lead, rest in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r0 = p.add_run()
        r0.text = '\u25b8  '
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.name = FONT
        r0.font.color.rgb = bullet_color
        if lead:
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.name = FONT
            r1.font.color.rgb = text_color
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.bold = False
        r2.font.name = FONT
        r2.font.color.rgb = TEXT_MUTED
    return box



def add_styled_table(slide, l, t, w, h, headers, rows, col_widths_in=None, accent=INDIGO):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(h))
    table = gtable.table

    if col_widths_in:
        for i, cw in enumerate(col_widths_in):
            table.columns[i].width = Inches(cw)

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.paragraphs[0].text = htext
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(12.5)
        run.font.bold = True
        run.font.name = FONT
        run.font.color.rgb = WHITE

    for r_idx, row in enumerate(rows, start=1):
        band = (r_idx % 2 == 0)
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_PANEL_ALT if band else BG_PANEL
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.paragraphs[0].text = str(val)
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(11.5)
            run.font.name = FONT
            run.font.color.rgb = TEXT_LIGHT
    return gtable


def add_chain_diagram(slide, top, steps):
    """steps: list of (line1, line2_or_None, fill_color). Glass-style rounded
    pill boxes connected by thin accent connector lines with an arrowhead."""
    n = len(steps)
    box_w, box_h, gap_w = 1.62, 0.62, 0.30
    total_w = n * box_w + (n - 1) * gap_w
    left = (SLIDE_W_IN - total_w) / 2
    x = left
    for i, (line1, line2, fill) in enumerate(steps):
        box = _rounded(slide, x, top, box_w, box_h, fill, radius=0.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = line1
        r1.font.size = Pt(11.5)
        r1.font.bold = True
        r1.font.name = FONT
        r1.font.color.rgb = WHITE
        if line2:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = line2
            r2.font.size = Pt(9.5)
            r2.font.name = FONT
            r2.font.color.rgb = RGBColor.from_string('E4E7FF')
        x += box_w
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 0.02), Inches(top + 0.18), Inches(gap_w - 0.04), Inches(0.26))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_MUTED
            arrow.line.fill.background()
            arrow.shadow.inherit = False
            x += gap_w


def add_screenshot_panel(slide, l, t, w, h, header_text, header_fill, image_file=None):
    """Glass panel with a small pill header chip, and either the real
    screenshot or a text placeholder inside."""
    add_glass_panel(slide, l, t, w, h)
    header_w = min(w - 0.2, max(2.4, len(header_text) * 0.085))
    hdr = _rounded(slide, l + 0.15, t + 0.12, header_w, 0.34, header_fill, radius=0.5)
    tf = hdr.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = header_text
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = WHITE

    img_top = t + 0.58
    img_h = h - 0.7
    img_l = l + 0.15
    img_w = w - 0.3

    if image_file:
        img_path = os.path.join(SCREEN_DIR, image_file)
        if os.path.exists(img_path):
            from PIL import Image
            im = Image.open(img_path)
            ratio = im.width / im.height
            fit_w = img_w
            fit_h = fit_w / ratio
            if fit_h > img_h:
                fit_h = img_h
                fit_w = fit_h * ratio
            pic_l = img_l + (img_w - fit_w) / 2
            pic_t = img_top + (img_h - fit_h) / 2
            slide.shapes.add_picture(img_path, Inches(pic_l), Inches(pic_t), Inches(fit_w), Inches(fit_h))
            return

    box = slide.shapes.add_textbox(Inches(img_l), Inches(img_top), Inches(img_w), Inches(img_h))
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = '(Insert Screenshot Here)'
    r2.font.size = Pt(12)
    r2.font.italic = True
    r2.font.name = FONT
    r2.font.color.rgb = TEXT_MUTED


def add_pill_header(slide, l, t, w, text, fill):
    hdr = _rounded(slide, l, t, w, 0.4, fill, radius=0.4)
    tf = hdr.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = WHITE
    return hdr



# =====================================================================
# SLIDE 1 - TITLE SLIDE (dark hero background, indigo glow accent bar,
# glass "presented by" card - same layout role as before, restyled dark)
# =====================================================================
title_slide = prs.slides.add_slide(blank_layout)
_rect(title_slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, BG_DARK)
_rect(title_slide, 0, 0, SLIDE_W_IN, 0.09, INDIGO)
_rect(title_slide, 0, 4.55, SLIDE_W_IN, 0.03, GREEN)

_textbox(title_slide, 0.6, 1.7, 12.1, 0.353, [('M.TECH MINI PROJECT', 15, True, AMBER)])
_textbox(title_slide, 0.6, 2.1, 12.1, 1.55,
         [[('FinGenius AI', 40, True, TEXT_LIGHT)],
          [('Smart Personal Finance Tracker Powered by Gemini AI', 22, True, RGBColor.from_string('C7CCF5'))]],
         line_spacing=1.1)
_textbox(title_slide, 0.6, 3.65, 12.1, 0.5,
         [('MERN Stack  \u2022  Budgets  \u2022  Goals  \u2022  Analytics  \u2022  AI Chat  \u2022  Receipt OCR  \u2022  Voice Entry', 15, False, TEXT_MUTED)])

add_glass_panel(title_slide, 2.4, 5.05, 8.5, 1.95, fill=BG_PANEL)
_textbox(title_slide, 2.75, 5.2, 7.8, 0.4, [('Presented by', 12, True, AMBER)])
_textbox(title_slide, 2.75, 5.52, 7.8, 0.42, [('[Student Name]  |  Reg. No. [Register Number]', 19, True, TEXT_LIGHT)])
_textbox(title_slide, 2.75, 5.98, 7.8, 0.337, [('M.Tech, Computer Science Engineering', 14, False, TEXT_MUTED)])
_textbox(title_slide, 2.75, 6.3, 7.8, 0.32, [('Department of Computer Science and Engineering', 13, False, TEXT_MUTED)])
_textbox(title_slide, 2.75, 6.6, 7.8, 0.34, [('[College Name]', 13, False, TEXT_MUTED)])



# =====================================================================
# SLIDE 2 - MOTIVATION: Problem Statement & Objectives (two glass panels)
# =====================================================================
s2 = add_content_slide('MOTIVATION', 'Problem Statement & Objectives', accent=INDIGO)
add_glass_panel(s2, 0.55, 1.55, 6.15, 5.25)
_textbox(s2, 0.85, 1.72, 5.6, 0.4, [('The Problem', 17, True, AMBER)])
add_bullets(s2, 0.85, 2.2, 5.6, 4.5, [
    ('Scattered records: ', 'people track expenses in notebooks or spreadsheets, with no single place for income, budget, and goals.'),
    ('No intelligent insight: ', 'free tools simply store numbers without explaining spending patterns or suggesting improvements.'),
    ('Manual entry is tedious: ', 'typing every transaction by hand causes many users to abandon tracking within weeks.'),
    ('No natural language help: ', 'users cannot simply ask "How much did I spend on food?" and get a direct answer.'),
    ('Privacy concerns: ', 'many finance apps need access to bank SMS or account details to work.'),
], bullet_color=AMBER, size=13)

add_glass_panel(s2, 6.9, 1.55, 5.9, 5.25)
_textbox(s2, 7.2, 1.72, 5.3, 0.4, [('Project Objectives', 17, True, GREEN)])
add_bullets(s2, 7.2, 2.2, 5.35, 4.5, [
    ('Unified tracking: ', 'record expenses and income and organise them into clear categories.'),
    ('Budget and goal planning: ', 'set monthly category budgets and savings goals with progress tracking.'),
    ('AI powered chat: ', 'integrate Google Gemini AI to answer questions using the user\'s own data.'),
    ('Faster data entry: ', 'support receipt scanning (OCR) and voice based expense entry.'),
    ('Visual analytics: ', 'present spending and saving trends through interactive charts.'),
], bullet_color=GREEN, size=13)



# =====================================================================
# SLIDE 3 - SYSTEM ARCHITECTURE (frontend / backend glass panels)
# =====================================================================
s3 = add_content_slide('SYSTEM DESIGN', 'System Architecture - Frontend & Backend', accent=GREEN)

add_pill_header(s3, 0.5, 1.5, 6.05, 'CLIENT SIDE  \u2022  React.js Frontend', INDIGO)
add_glass_panel(s3, 0.5, 2.0, 6.05, 4.95)
add_bullets(s3, 0.85, 2.3, 5.4, 4.5, [
    ('React 18 + Vite: ', 'component based single page application'),
    ('Tailwind CSS: ', 'responsive glassmorphism styling, dark/light theme'),
    ('Framer Motion: ', 'smooth page and card animations'),
    ('Recharts: ', 'interactive analytics charts'),
    ('React Router v6: ', 'client side routing between pages'),
    ('Axios: ', 'REST API communication with the backend'),
    ('Pages: ', 'Landing, Login, Register, Dashboard, Expenses, Income, Budget, Goals, Analytics, Chat, Settings'),
], bullet_color=INDIGO, size=12.5)

add_pill_header(s3, 6.75, 1.5, 6.05, 'SERVER SIDE  \u2022  Node.js Backend', GREEN)
add_glass_panel(s3, 6.75, 2.0, 6.05, 4.95)
add_bullets(s3, 7.1, 2.3, 5.4, 4.5, [
    ('Node.js + Express.js: ', 'REST API server with route-controller structure'),
    ('MongoDB Atlas + Mongoose: ', 'cloud database and schema modelling'),
    ('JWT + bcryptjs: ', 'token based authentication, hashed passwords'),
    ('Google Gemini AI: ', 'AI chat assistant with financial context'),
    ('Tesseract.js: ', 'OCR engine for receipt scanning'),
    ('Multer: ', 'file upload handling for receipt images'),
    ('Middleware: ', 'JWT auth guard, rate limiter, input validator, error handler'),
], bullet_color=GREEN, size=12.5)



# =====================================================================
# SLIDE 4 - HOW IT WORKS: Methodology (glass chain diagram + explanation)
# =====================================================================
s4 = add_content_slide('HOW IT WORKS', 'Methodology - From Login to AI Insight', accent=AMBER)

add_chain_diagram(s4, 1.55, [
    ('Register &\nLogin', 'JWT + bcrypt', INDIGO),
    ('Dashboard', 'Balance overview', GREEN),
    ('Add\nTransaction', 'Manual/OCR/Voice', AMBER),
    ('Budget &\nGoals', 'Category limits', INDIGO),
    ('Analytics', 'Recharts trends', GREEN),
    ('AI Chat', 'Gemini financial Q&A', AMBER),
])

add_bullets(s4, 0.55, 2.55, 12.3, 4.4, [
    ('User Authentication: ', 'a user registers with name, email, and password; passwords are hashed with bcryptjs (12 salt rounds) and login issues a 7-day JWT token used to authorise every later request.'),
    ('Recording Transactions: ', 'expenses and income can be entered manually, extracted from a receipt photo using Tesseract.js OCR, or logged by speaking through the Web Speech API.'),
    ('Budget and Goal Tracking: ', 'users set category wise monthly budgets and savings goals with a target amount; the system tracks actual spending and contribution progress automatically.'),
    ('Analytics Generation: ', 'Recharts renders monthly expense trend, income vs expense, category breakdown, and savings trend charts from the stored MongoDB data.'),
    ('AI Powered Chat: ', 'the backend sends the user\u2019s financial summary (expenses by category, income, budget status, goals, recent transactions) to Google Gemini 1.5 Flash, which answers questions in plain language.'),
], bullet_color=AMBER, size=14)



# =====================================================================
# SLIDE 5 - TECHNOLOGY STACK (table + dashboard screenshot panel)
# =====================================================================
s5 = add_content_slide('IMPLEMENTATION', 'Technology Stack', accent=INDIGO)

_textbox(s5, 0.55, 1.5, 6.0, 0.35, [('Frontend and Backend Technologies', 14, True, TEXT_LIGHT)])
add_styled_table(s5, 0.55, 1.9, 6.3, 3.4,
                  ['Technology', 'Purpose'],
                  [
                      ['React 18 + Vite', 'Frontend UI'],
                      ['Tailwind CSS', 'Responsive styling'],
                      ['Node.js + Express.js', 'Backend REST APIs'],
                      ['MongoDB + Mongoose', 'Database & schema'],
                      ['JWT + bcryptjs', 'Auth & security'],
                      ['Google Gemini AI', 'AI chat assistant'],
                      ['Tesseract.js', 'Receipt OCR'],
                      ['Web Speech API', 'Voice entry'],
                      ['Recharts', 'Analytics charts'],
                  ],
                  col_widths_in=[3.0, 3.3], accent=INDIGO)

add_screenshot_panel(s5, 7.1, 1.55, 5.7, 4.95, 'Dashboard \u2022 Financial Overview', INDIGO, 'Screenshot 2026-07-13 205528.png')



# =====================================================================
# SLIDE 6 - LIVE SYSTEM: Application Walkthrough (two screenshot panels)
# =====================================================================
s6 = add_content_slide('LIVE SYSTEM', 'Application Walkthrough', accent=GREEN)

add_screenshot_panel(s6, 0.5, 1.55, 5.9, 4.95, 'Expenses \u2022 Search \u2022 Filter \u2022 Scan \u2022 Voice', GREEN, 'Screenshot 2026-07-13 205600.png')
add_screenshot_panel(s6, 6.9, 1.55, 5.9, 4.95, 'AI Chat \u2022 Gemini Powered Assistant', AMBER, 'Screenshot 2026-07-13 205745.png')



# =====================================================================
# SLIDE 7 - COMPARISON WITH EXISTING SYSTEMS (table + merits/limitations)
# =====================================================================
s7 = add_content_slide('COMPARISON', 'Comparison with Existing Systems', accent=AMBER)

_textbox(s7, 0.55, 1.5, 8.0, 0.35, [('Existing Methods vs. FinGenius AI', 14, True, TEXT_LIGHT)])
add_styled_table(s7, 0.55, 1.9, 12.25, 2.1,
                  ['Feature', 'Manual / Excel', 'Typical Finance Apps', 'FinGenius AI'],
                  [
                      ['Security', 'None / basic', 'Varies by app', 'JWT + bcrypt hashing'],
                      ['AI Assistance', 'Not available', 'Rare or basic', 'Gemini AI chatbot'],
                      ['Expense Entry', 'Manual only', 'Manual entry', 'Manual, OCR, and voice'],
                      ['Cost', 'Free', 'Free to paid', 'Free for this project'],
                  ],
                  col_widths_in=[2.6, 2.9, 3.35, 3.4], accent=AMBER)

add_pill_header(s7, 0.55, 4.35, 6.05, 'Merits', GREEN)
add_glass_panel(s7, 0.55, 4.85, 6.05, 2.05)
add_bullets(s7, 0.85, 5.0, 5.5, 1.8, [
    (None, 'AI powered insights based on the user\u2019s real financial data'),
    (None, 'Multiple ways to log expenses - manual, receipt scan, voice'),
    (None, 'Completely free with all features unlocked'),
    (None, 'Secure - JWT authentication and hashed passwords'),
], bullet_color=GREEN, size=12)

add_pill_header(s7, 6.85, 4.35, 5.95, 'Limitations', AMBER)
add_glass_panel(s7, 6.85, 4.85, 5.95, 2.05)
add_bullets(s7, 7.15, 5.0, 5.45, 1.8, [
    (None, 'No direct bank account integration yet - entries are manual or scanned'),
    (None, 'Requires an internet connection for all operations'),
    (None, 'OCR accuracy depends on the quality of the receipt image'),
    (None, 'Voice recognition works best in Chrome or Edge browsers'),
], bullet_color=AMBER, size=12)



# =====================================================================
# SLIDE 8 - WRAP-UP: Conclusion & Future Work
# =====================================================================
s8 = add_content_slide('WRAP-UP', 'Conclusion & Future Work', accent=INDIGO)

add_pill_header(s8, 0.55, 1.5, 6.05, 'Conclusion', INDIGO)
add_glass_panel(s8, 0.55, 2.0, 6.05, 4.3)
add_bullets(s8, 0.85, 2.2, 5.5, 4.0, [
    (None, 'Delivered a full-stack personal finance tracker combining expense, income, budget, and goal management with an AI chat assistant.'),
    (None, 'Integrated Google Gemini AI so users get personalised, data based answers instead of generic financial tips.'),
    (None, 'Added Receipt OCR and Voice Entry to reduce the effort of manual data entry.'),
    (None, 'Implemented secure JWT authentication and a fully responsive, modern user interface.'),
    (None, 'Provides a practical, free, all-in-one alternative to scattered manual and spreadsheet based tracking.'),
], bullet_color=INDIGO, size=13)

add_pill_header(s8, 6.85, 1.5, 5.95, 'Future Work', GREEN)
add_glass_panel(s8, 6.85, 2.0, 5.95, 4.3)
add_bullets(s8, 7.15, 2.2, 5.45, 4.0, [
    (None, 'Bank account and UPI integration for automatic transaction import'),
    (None, 'Improving OCR accuracy for a wider variety of receipt formats'),
    (None, 'Multi-language support for the UI and voice entry'),
    (None, 'Investment tracking alongside income and expenses'),
    (None, 'Predictive expense forecasting using past spending patterns'),
    (None, 'A dedicated mobile application with push notifications'),
], bullet_color=GREEN, size=13)

_rect(s8, 0, 6.35, SLIDE_W_IN, 0.75, INDIGO)
_textbox(s8, 0, 6.35, SLIDE_W_IN, 0.75, [('Thank You', 20, True, WHITE)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUTPUT)
print(f'Presentation saved to {OUTPUT}')
